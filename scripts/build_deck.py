#!/usr/bin/env python3
"""Final Project STI Kel 5 - presentation deck builder (python-pptx)."""
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PM = os.path.join(HERE, '.pmtmp')
PREV = os.path.join(HERE, 'previews')
OUT = os.path.join(HERE, 'Final-Project-STI-Kel5-Presentation.pptx')

INK   = RGBColor(0x1C, 0x6F, 0x8C)   # deep teal
DARK  = RGBColor(0x12, 0x3A, 0x49)
ACC   = RGBColor(0xD9, 0x8A, 0x4B)   # terracotta
GRN   = RGBColor(0x2E, 0x8B, 0x6F)
LT    = RGBColor(0x9E, 0xCA, 0xE1)
BG    = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF1, 0xF6, 0xF8)
GREY  = RGBColor(0x5A, 0x6A, 0x70)
INKL  = RGBColor(0xE2, 0xEE, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = 'Segoe UI'
FONTH = 'Segoe UI Semibold'

prs = Presentation()
prs.slide_width  = IN(13.333)
prs.slide_height = IN(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    return s

def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line=1.05):
    tb = s.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs[0], tuple): runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        p.line_spacing = line
        for (t, sz, col, bold, *rest) in para:
            r = p.add_run(); r.text = t; f = r.font
            f.size = Pt(sz); f.color.rgb = col; f.bold = bold
            f.name = FONTH if bold else FONT
            if rest and rest[0]: f.italic = True
    return tb

def bullet(s, x, y, w, h, items, sz=15, col=DARK, gap=7, mk=INK, line=1.1):
    tb = s.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h)); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        lvl = 0; txt = it
        if isinstance(it, tuple): lvl, txt = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = line
        r = p.add_run(); r.text = ('   ' if lvl else '') + ('–  ' if lvl else '•  ')
        r.font.size = Pt(sz); r.font.color.rgb = mk if not lvl else GREY; r.font.bold = False; r.font.name = FONT
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz if not lvl else sz-1); r.font.color.rgb = col if not lvl else GREY
        r.font.name = FONT
    return tb

def header(s, kicker, title, n):
    box(s, 0, 0, 13.333, 1.28, fill=BG)
    box(s, 0.6, 0.46, 0.14, 0.62, fill=INK)
    text(s, 0.92, 0.40, 11, 0.35, [[(kicker.upper(), 11.5, ACC, True)]], sp_after=0)
    text(s, 0.92, 0.62, 11.4, 0.6, [[(title, 25, INK, True)]], sp_after=0)
    box(s, 0.6, 1.30, 12.133, 0.018, fill=INKL)
    text(s, 12.2, 6.96, 1.0, 0.4, [[(f'{n:02d}', 11, GREY, False)]], align=PP_ALIGN.RIGHT, sp_after=0)
    text(s, 0.6, 6.96, 6, 0.4, [[('Final Project STI · Kelompok 5', 10, GREY, False)]], sp_after=0)

def pill(s, x, y, w, label, color):
    h = 0.34
    p = box(s, x, y, w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p.adjustments[0] = 0.5
    text(s, x, y+0.02, w, h, [[(label, 11, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp_after=0)

n = 0
# ---------- 1. COVER ----------
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=INK)
box(s, 0, 0, 13.333, 7.5, fill=INK)
# decorative hex-ish panels
box(s, 8.7, -1.2, 6, 6, fill=DARK, shape=MSO_SHAPE.DIAMOND)
box(s, 10.8, 4.2, 5, 5, fill=RGBColor(0x21,0x55,0x66), shape=MSO_SHAPE.OVAL)
box(s, 0, 6.95, 13.333, 0.55, fill=DARK)
text(s, 0.95, 1.15, 11, 0.4, [[('MAGISTER MANAJEMEN TEKNOLOGI · INSTITUT TEKNOLOGI SEPULUH NOPEMBER', 12, LT, True)]], sp_after=0)
text(s, 0.9, 1.95, 11.4, 2.0,
     [[('Final Project', 50, WHITE, True)],
      [('Sistem & Teknologi Informasi', 30, LT, False)]], sp_after=8, line=1.0)
box(s, 0.95, 3.95, 3.2, 0.05, fill=ACC)
text(s, 0.95, 4.2, 11, 0.5, [[('Analisis & Monitoring Data dengan Pendekatan CRISP-DM', 17, WHITE, False)]], sp_after=0)
text(s, 0.95, 4.75, 11, 0.4, [[('Dua studi kasus · Dashboard interaktif · Manajemen proyek GitHub', 13, LT, False)]], sp_after=0)
# member block
box(s, 0.95, 5.5, 6.6, 1.25, fill=DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 1.2, 5.62, 6.2, 1.1,
     [[('KELOMPOK 5', 11, ACC, True)],
      [('Moch Chesa Nur Hidayat', 13.5, WHITE, False), ('   ·   [Anggota 2]   ·   [Anggota 3]', 13.5, LT, False)],
      [('GitHub: sianida26, Misbahulmunir26, …   |   Dosen: [isi nama dosen]', 10.5, LT, False)]],
     sp_after=3, line=1.05)
text(s, 8.1, 5.95, 4.3, 0.5, [[('23 Juni 2026', 13, WHITE, True)]], align=PP_ALIGN.RIGHT, sp_after=0)
n += 1

# ---------- 2. AGENDA ----------
n += 1; s = slide(); header(s, 'Outline', 'Agenda Presentasi', n)
items = [
    ('01', 'Latar Belakang & Tujuan', 'Konteks dua studi kasus dan business questions'),
    ('02', 'Metodologi CRISP-DM', 'Kerangka kerja pengolahan data end-to-end'),
    ('03', 'Manajemen Proyek', 'Pelacakan pekerjaan via GitHub Projects (kanban)'),
    ('04', 'Studi Kasus 1: Inventory CGS-1', 'Monitoring kesiapan & kalibrasi instrumen'),
    ('05', 'Studi Kasus 2: Penjualan Ritel', 'Tren, pendorong revenue, musiman, diskon'),
    ('06', 'Kesimpulan & Rekomendasi', 'Temuan utama dan tindak lanjut'),
]
y = 1.62
for num, t, d in items:
    box(s, 0.9, y, 0.95, 0.82, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 0.9, y+0.02, 0.95, 0.8, [[(num, 23, INK, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp_after=0)
    text(s, 2.05, y+0.10, 10.4, 0.8,
         [[(t, 17, DARK, True)], [(d, 12.5, GREY, False)]], sp_after=2, line=1.0)
    y += 0.9
n_agenda = n

# ---------- 3. LATAR BELAKANG ----------
n += 1; s = slide(); header(s, 'Pendahuluan', 'Latar Belakang & Tujuan', n)
text(s, 0.9, 1.55, 11.6, 0.7,
     [[('Tim mengangkat ', 14.5, DARK, False), ('dua studi kasus pengolahan data nyata', 14.5, INK, True),
       (' dari lingkungan kerja anggota, lalu mengolahnya secara terstruktur (CRISP-DM) hingga menjadi dashboard pemantauan.', 14.5, DARK, False)]],
     sp_after=0, line=1.15)
# two case cards
cards = [
    (0.9, 'STUDI KASUS 1', 'Inventory Instrumen CGS-1', INK,
     ['Monitoring kesiapan instrumen lapangan (transmitter, valve, actuator) pada Custody Gas Station.',
      'Fokus: status aset, ketergantungan vendor, dan gap data kalibrasi (keselamatan & kepatuhan).',
      'Output: dashboard Looker Studio.']),
    (7.0, 'STUDI KASUS 2', 'Kinerja Penjualan Ritel Online', GRN,
     ['Data penjualan B2C agregat mingguan, 2022–2026 (~3,5 tahun).',
      'Fokus: tren & pertumbuhan, pendorong revenue, pola musiman, efektivitas diskon.',
      'Output: dashboard interaktif Plotly.']),
]
for x, k, t, c, pts in cards:
    box(s, x, 2.45, 5.45, 4.2, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, x, 2.45, 5.45, 0.12, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.35, 2.72, 4.8, 0.3, [[(k, 11, c, True)]], sp_after=0)
    text(s, x+0.35, 3.02, 4.8, 0.7, [[(t, 19, DARK, True)]], sp_after=0, line=1.0)
    bullet(s, x+0.35, 3.9, 4.8, 2.6, pts, sz=13, gap=9, mk=c)
n_case_intro = n

# ---------- 4. CRISP-DM ----------
n += 1; s = slide(); header(s, 'Metodologi', 'Kerangka Kerja CRISP-DM', n)
text(s, 0.9, 1.5, 11.6, 0.5,
     [[('Setiap studi kasus mengikuti enam fase CRISP-DM secara berurutan, dari pemahaman bisnis hingga deployment dashboard.', 14, GREY, False)]], sp_after=0, line=1.1)
phases = [
    ('1', 'Business\nUnderstanding', 'Tujuan & business questions', INK),
    ('2', 'Data\nUnderstanding', 'Load data, profil kualitas', INK),
    ('3', 'Data\nPreparation', 'Cleaning, kolom turunan (KPI)', INK),
    ('4', 'Modeling /\nAnalisis', 'Agregasi & visualisasi', GRN),
    ('5', 'Evaluation', 'Temuan & rekomendasi', GRN),
    ('6', 'Deployment', 'Dashboard monitoring', ACC),
]
x = 0.9; w = 1.86; gap = 0.13; y = 2.45
for i,(num,t,d,c) in enumerate(phases):
    box(s, x, y, w, 2.5, fill=WHITE, line=INKL, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    circ = box(s, x+w/2-0.33, y+0.28, 0.66, 0.66, fill=c, shape=MSO_SHAPE.OVAL)
    text(s, x+w/2-0.33, y+0.30, 0.66, 0.66, [[(num, 22, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp_after=0)
    text(s, x+0.1, y+1.08, w-0.2, 0.85, [[(ln, 13.5, DARK, True)] for ln in t.split('\n')], align=PP_ALIGN.CENTER, sp_after=0, line=1.0)
    text(s, x+0.12, y+1.95, w-0.24, 0.5, [[(d, 10.5, GREY, False)]], align=PP_ALIGN.CENTER, sp_after=0, line=1.0)
    if i < 5:
        box(s, x+w+0.005, y+1.05, gap+0.02, 0.12, fill=LT, shape=MSO_SHAPE.CHEVRON)
    x += w + gap
text(s, 0.9, 5.25, 11.6, 0.5,
     [[('Pendekatan ini memastikan analisis selalu berangkat dari pertanyaan bisnis dan diakhiri dengan artefak yang dapat dipantau, bukan sekadar grafik lepas.', 13, GREY, False, True)]], sp_after=0, line=1.15)
n_crisp = n

# ---------- 5. PM intro ----------
n += 1; s = slide(); header(s, 'Manajemen Proyek', 'Pelacakan Pekerjaan dengan GitHub Projects', n)
text(s, 0.9, 1.5, 11.6, 0.8,
     [[('Pekerjaan kelompok dikelola pada ', 14.5, DARK, False), ('board kanban GitHub Projects', 14.5, INK, True),
       (', dengan setiap fase CRISP-DM dari kedua studi kasus dipecah menjadi kartu tugas dan dilacak melalui kolom ', 14.5, DARK, False),
       ('Todo → In Progress → Done', 14.5, GRN, True), ('.', 14.5, DARK, False)]], sp_after=0, line=1.15)
why = [
    ('Transparansi', 'Status tiap pekerjaan terlihat satu layar, tidak ada tugas yang “hilang”.'),
    ('Akuntabilitas', 'Pemecahan tugas yang jelas memudahkan pembagian kerja antar anggota.'),
    ('Keterlacakan', 'Progres terukur (12 selesai, 2 berjalan, 2 antre) dan terhubung ke repositori kode.'),
]
x = 0.9
for t, d in why:
    box(s, x, 2.65, 3.78, 1.7, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, x+0.32, 2.95, 0.5, 0.08, fill=ACC)
    text(s, x+0.32, 3.12, 3.2, 0.4, [[(t, 16, INK, True)]], sp_after=0)
    text(s, x+0.32, 3.55, 3.2, 0.9, [[(d, 12, GREY, False)]], sp_after=0, line=1.1)
    x += 3.97
# repo / board links
box(s, 0.9, 4.7, 11.53, 1.55, fill=DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 1.25, 4.9, 11, 0.4, [[('TAUTAN PROYEK', 11, ACC, True)]], sp_after=0)
text(s, 1.25, 5.25, 11, 0.45, [[('Repository : ', 14, LT, True), ('github.com/sianida26/sti-final-project-kel5', 14, WHITE, False)]], sp_after=0)
text(s, 1.25, 5.68, 11, 0.45, [[('Board       : ', 14, LT, True), ('github.com/users/sianida26/projects/2', 14, WHITE, False)]], sp_after=0)
n_pm = n

# ---------- 6. PM board screenshot ----------
n += 1; s = slide(); header(s, 'Manajemen Proyek', 'Board Kanban: Status Pekerjaan', n)
# board image left
img = os.path.join(PM, 'board.png')
from PIL import Image as PImage
iw, ih = PImage.open(img).size
disp_w = 8.0; disp_h = disp_w * ih / iw
if disp_h > 5.15: disp_h = 5.15; disp_w = disp_h * iw / ih
box(s, 0.85, 1.5, disp_w+0.1, disp_h+0.1, fill=None, line=INKL, line_w=1.2)
s.shapes.add_picture(img, IN(0.9), IN(1.55), IN(disp_w), IN(disp_h))
# right summary
rx = 9.3
text(s, rx, 1.7, 3.4, 0.4, [[('RINGKASAN STATUS', 11, ACC, True)]], sp_after=0)
stats = [('12', 'Done', INK), ('2', 'In Progress', ACC), ('2', 'Todo', GREY)]
yy = 2.15
for v, lb, c in stats:
    box(s, rx, yy, 3.45, 0.92, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, rx, yy, 0.12, 0.92, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, rx+0.35, yy+0.04, 1.1, 0.85, [[(v, 30, c, True)]], anchor=MSO_ANCHOR.MIDDLE, sp_after=0)
    text(s, rx+1.5, yy+0.04, 1.9, 0.85, [[(lb, 15, DARK, True)]], anchor=MSO_ANCHOR.MIDDLE, sp_after=0)
    yy += 1.06
text(s, rx, yy+0.1, 3.45, 1.6,
     [[('16 kartu total', 13, INK, True)],
      [('Kartu = fase CRISP-DM kedua studi kasus + setup, dokumentasi, dan pengumpulan.', 11.5, GREY, False)]],
     sp_after=6, line=1.15)
n_board = n

# ---------- 7. WBS mapping ----------
n += 1; s = slide(); header(s, 'Manajemen Proyek', 'Pemetaan Tugas ke Fase CRISP-DM (WBS)', n)
rows = [
    ('Fase CRISP-DM', 'Inventory CGS-1', 'Penjualan Ritel', True),
    ('Business Understanding', 'Definisi 3 business question', 'Definisi 4 business question', False),
    ('Data Understanding', 'Load Excel, profil missing', 'Ingest data mingguan, kamus data', False),
    ('Data Preparation', 'Cleaning, flag kalibrasi/serial', 'KPI turunan: AOV, MA, disc%', False),
    ('Modeling / Analisis', 'Status, vendor, coverage kalibrasi', 'Tren, musiman, AOV, diskon', False),
    ('Evaluation', 'Gap kalibrasi, watchlist aset', 'Pertumbuhan, over-discounting', False),
    ('Deployment', 'Dashboard Looker Studio', 'Dashboard Plotly interaktif', False),
]
x0, y0 = 0.9, 1.7; cw = [4.0, 3.85, 3.78]; rh = 0.72
yy = y0
for ri, (a,b,c,hd) in enumerate(rows):
    xx = x0
    for ci, val in enumerate((a,b,c)):
        fill = INK if hd else (PANEL if ri%2 else WHITE)
        cell = box(s, xx, yy, cw[ci], rh, fill=fill, line=INKL, line_w=0.75)
        col = WHITE if hd else (INK if ci==0 else DARK)
        bold = hd or ci==0
        text(s, xx+0.18, yy, cw[ci]-0.3, rh, [[(val, 12.5 if not hd else 13, col, bold)]],
             anchor=MSO_ANCHOR.MIDDLE, sp_after=0, line=1.0)
        xx += cw[ci]
    yy += rh
n_wbs = n

# ---------- 8. Inventory BU ----------
n += 1; s = slide(); header(s, 'Studi Kasus 1 · Inventory CGS-1', 'Konteks & Business Questions', n)
text(s, 0.9, 1.55, 11.6, 1.0,
     [[('CGS-1 (Custody Gas Station)', 14.5, INK, True),
       (' memiliki sejumlah instrumen lapangan yang menopang keandalan, keselamatan, dan kepatuhan operasi. Data instrumen tersedia namun belum diolah dan kualitasnya belum dievaluasi.', 14.5, DARK, False)]],
     sp_after=0, line=1.15)
box(s, 0.9, 2.75, 11.53, 0.05, fill=INKL)
qs = [
    ('Q1', 'Komposisi & status', 'Bagaimana sebaran instrumen aktif vs tidak aktif?'),
    ('Q2', 'Jenis & vendor', 'Equipment class dan manufacturer apa yang dominan (ketergantungan vendor)?'),
    ('Q3', 'Kelengkapan kalibrasi', 'Instrumen mana yang belum punya data kalibrasi (potensi gap keselamatan/kepatuhan)?'),
]
yy = 3.1
for k, t, d in qs:
    box(s, 0.9, yy, 1.0, 1.0, fill=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 0.9, yy+0.02, 1.0, 1.0, [[(k, 22, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp_after=0)
    text(s, 2.1, yy+0.12, 10.2, 0.9, [[(t, 16.5, DARK, True)], [(d, 13, GREY, False)]], sp_after=2, line=1.05)
    yy += 1.18
n_inv_bu = n

# ---------- 9. Inventory findings ----------
n += 1; s = slide(); header(s, 'Studi Kasus 1 · Inventory CGS-1', 'Temuan & Rekomendasi', n)
text(s, 0.9, 1.5, 5.7, 0.4, [[('TEMUAN UTAMA', 12, ACC, True)]], sp_after=0)
bullet(s, 0.9, 1.95, 5.7, 4.6, [
    'Status: mayoritas instrumen Active (~94%), ketersediaan aset baik.',
    'Komposisi didominasi Actuator & Pressure Transmitter; vendor terkonsentrasi pada Rosemount & Honeywell → ada ketergantungan vendor.',
    'Insight utama: gap data kalibrasi signifikan; sebagian instrumen tidak memiliki Calibrate Range tercatat.',
    'Instrumen Active tanpa data kalibrasi menjadi watchlist prioritas (risiko kepatuhan/keselamatan).',
    'Sebagian Serial Number kosong → menyulitkan keterlacakan aset & klaim garansi.',
], sz=13.5, gap=11, mk=INK)
# right reco panel
box(s, 7.0, 1.5, 5.43, 5.05, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 7.35, 1.75, 4.8, 0.4, [[('REKOMENDASI', 12, GRN, True)]], sp_after=0)
bullet(s, 7.35, 2.2, 4.8, 3.0, [
    'Lengkapi Calibrate Range & Serial Number untuk instrumen prioritas (Active tanpa kalibrasi).',
    'Jadikan % cakupan kalibrasi sebagai KPI monitoring berkala di dashboard.',
    'Standardisasi input data di sumber untuk menghindari field kosong.',
], sz=13, gap=10, mk=GRN)
box(s, 7.35, 5.25, 4.7, 1.05, fill=WHITE, line=INKL, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 7.6, 5.38, 4.3, 0.85,
     [[('Keterbatasan: ', 11.5, ACC, True), ('sampel kecil (18 instrumen), satu area, snapshot satu waktu → analisis deskriptif, belum prediktif.', 11.5, GREY, False)]],
     sp_after=0, line=1.1)
n_inv_find = n

# ---------- 10. Sales BU ----------
n += 1; s = slide(); header(s, 'Studi Kasus 2 · Penjualan Ritel', 'Konteks & Business Questions', n)
text(s, 0.9, 1.55, 11.6, 0.9,
     [[('Sebuah retailer online B2C', 14.5, GRN, True),
       (' ingin memahami kinerja penjualannya selama ~3,5 tahun (data agregat mingguan, anonim) untuk mendukung keputusan stok, promo, dan target.', 14.5, DARK, False)]],
     sp_after=0, line=1.15)
qs = [
    ('Q1', 'Tren & pertumbuhan', 'Bagaimana tren revenue, transaksi, dan pelanggan dari waktu ke waktu?'),
    ('Q2', 'Pendorong revenue', 'Pertumbuhan ditarik oleh volume transaksi atau nilai per transaksi (AOV)?'),
    ('Q3', 'Pola musiman', 'Adakah bulan ramai/sepi untuk perencanaan stok & campaign?'),
    ('Q4', 'Efektivitas diskon', 'Seberapa efektif strategi diskon terhadap revenue?'),
]
yy = 2.75
for k, t, d in qs:
    box(s, 0.9, yy, 0.92, 0.86, fill=GRN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 0.9, yy+0.02, 0.92, 0.86, [[(k, 19, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp_after=0)
    text(s, 2.0, yy+0.06, 10.3, 0.8, [[(t, 16, DARK, True)], [(d, 12.5, GREY, False)]], sp_after=2, line=1.05)
    yy += 1.0
n_sales_bu = n

# ---------- 11. Sales dashboard ----------
n += 1; s = slide(); header(s, 'Studi Kasus 2 · Penjualan Ritel', 'Dashboard Interaktif (Plotly)', n)
img = os.path.join(PREV, 'plotly_dashboard_preview.png')
iw, ih = PImage.open(img).size
disp_h = 5.15; disp_w = disp_h * iw / ih
box(s, 0.85, 1.5, disp_w+0.1, disp_h+0.1, fill=None, line=INKL, line_w=1.2)
s.shapes.add_picture(img, IN(0.9), IN(1.55), IN(disp_w), IN(disp_h))
rx = 0.9 + disp_w + 0.45
text(s, rx, 1.7, 12.4-rx, 0.4, [[('KOMPONEN', 11, ACC, True)]], sp_after=0)
bullet(s, rx, 2.15, 12.4-rx, 4.4, [
    'Scorecard KPI: revenue, transaksi, AOV, rata-rata diskon.',
    'Tren revenue mingguan + rata-rata 4 minggu (range slider, dropdown tahun).',
    'Revenue per tahun & pola musiman per bulan.',
    'AOV per tahun & scatter diskon vs revenue.',
], sz=12.5, gap=10, mk=GRN)
text(s, rx, 5.5, 12.4-rx, 0.9,
     [[('Interaktif: ', 11.5, INK, True), ('hover, zoom, filter tahun, langsung di notebook.', 11.5, GREY, False)]],
     sp_after=0, line=1.1)
n_sales_dash = n

# ---------- 12. Sales findings ----------
n += 1; s = slide(); header(s, 'Studi Kasus 2 · Penjualan Ritel', 'Temuan & Rekomendasi', n)
kpis = [('5×', 'Pertumbuhan revenue\n2023 → 2025', GRN), ('AOV ↑', 'Pendorong utama\n(bukan volume)', INK),
        ('~1,3', 'Unit per transaksi\n(keranjang kecil)', ACC), ('Jun–Jul', 'Puncak musiman\n(juga Januari)', DARK)]
x = 0.9
for v, lb, c in kpis:
    box(s, x, 1.55, 2.78, 1.5, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, x, 1.55, 2.78, 0.1, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.2, 1.72, 2.4, 0.5, [[(v, 26, c, True)]], sp_after=0)
    text(s, x+0.2, 2.28, 2.5, 0.7, [[(ln, 11.5, GREY, False)] for ln in lb.split('\n')], sp_after=0, line=1.0)
    x += 2.92
text(s, 0.9, 3.35, 5.7, 0.4, [[('TEMUAN', 12, ACC, True)]], sp_after=0)
bullet(s, 0.9, 3.78, 5.75, 3.0, [
    'Revenue tahunan naik ~5× dalam dua tahun (2023 → 2025).',
    'Pendorong = nilai (AOV naik tajam), bukan sekadar jumlah transaksi.',
    'Diskon belum efektif: rata-rata ~36% namun korelasi diskon-revenue ~0 → indikasi over-discounting.',
], sz=12.5, gap=9, mk=INK)
box(s, 7.0, 3.35, 5.43, 3.2, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 7.35, 3.55, 4.8, 0.4, [[('REKOMENDASI', 12, GRN, True)]], sp_after=0)
bullet(s, 7.35, 4.0, 4.8, 2.5, [
    'Pertahankan momentum AOV via bundling/upsell (keranjang masih 1,3 item).',
    'Evaluasi strategi diskon: uji kurangi diskon, ukur dampak ke revenue.',
    'Manfaatkan puncak musiman (Jun–Jul, Jan) untuk stok & promo.',
    'Dorong akuisisi/retensi pelanggan (jumlah pelanggan relatif stagnan).',
], sz=12, gap=8, mk=GRN)
n_sales_find = n

# ---------- 13. Kesimpulan ----------
n += 1; s = slide(); header(s, 'Penutup', 'Kesimpulan', n)
bullet(s, 0.9, 1.7, 11.5, 4.0, [
    'Pendekatan CRISP-DM berhasil mengubah dua dataset mentah menjadi artefak pemantauan yang menjawab pertanyaan bisnis nyata.',
    'Inventory CGS-1: aset siap (~94% aktif), namun gap data kalibrasi menjadi prioritas tindak lanjut keselamatan & kepatuhan.',
    'Penjualan ritel: pertumbuhan kuat (~5×) ditarik kenaikan AOV; strategi diskon perlu dievaluasi (indikasi over-discounting).',
    'Manajemen proyek via GitHub Projects menjaga pekerjaan kelompok transparan, terbagi, dan terlacak (12 selesai · 2 berjalan · 2 antre).',
    'Tindak lanjut terbuka: bersihkan outlier data penjualan dan finalisasi dashboard Looker Studio inventory.',
], sz=15, gap=16, mk=INK, line=1.15)
n_concl = n

# ---------- 14. Thanks ----------
n += 1; s = slide()
box(s, 0, 0, 13.333, 7.5, fill=INK)
box(s, 9.0, -1.5, 6.5, 6.5, fill=DARK, shape=MSO_SHAPE.DIAMOND)
box(s, 10.6, 4.5, 5, 5, fill=RGBColor(0x21,0x55,0x66), shape=MSO_SHAPE.OVAL)
text(s, 0.95, 2.7, 11, 1.2, [[('Terima Kasih', 46, WHITE, True)]], sp_after=0)
box(s, 1.0, 3.95, 3.0, 0.05, fill=ACC)
text(s, 0.98, 4.25, 11, 0.5, [[('Final Project STI · Kelompok 5 · MMT ITS', 15, LT, False)]], sp_after=0)
text(s, 0.98, 4.75, 11, 0.5, [[('github.com/sianida26/sti-final-project-kel5', 13, WHITE, False)]], sp_after=0)

prs.save(OUT)
print('Saved', OUT, '·', len(prs.slides.__iter__.__self__._sldIdLst), 'slides')
